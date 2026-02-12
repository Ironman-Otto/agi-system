import sys
from pathlib import Path
import shutil
import json

from docx import Document
import pdfplumber
from pptx import Presentation


def extract_docx(path: Path) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def extract_pdf(path: Path) -> str:
    text = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text.append(t)
    return "\n".join(text)


def extract_pptx(path: Path) -> str:
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)
    return "\n".join(slides_text)


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()

    try:
        if ext == ".docx":
            return extract_docx(path)

        elif ext == ".pdf":
            return extract_pdf(path)

        elif ext == ".pptx":
            try:
                return extract_pptx(path)
            except Exception as e:
                print(f"⚠ PPTX extraction failed: {path.name} — {e}")
                return ""

        elif ext in [".md", ".txt"]:
            return path.read_text(encoding="utf-8", errors="ignore")

        else:
            return ""

    except Exception as e:
        print(f"⚠ Extraction error: {path.name} — {e}")
        return ""


def main(raw_folder: str, output_folder: str):

    raw_path = Path(raw_folder)
    out_path = Path(output_folder)

    out_path.mkdir(parents=True, exist_ok=True)

    index = []
    failed_files = []

    for file in raw_path.rglob("*"):

        if file.is_dir():
            continue

        if "diagrams" in file.parts:
            continue

        print(f"Processing: {file.name}")

        text = extract_text(file)

        if not text.strip():
            failed_files.append(str(file))
            continue

        output_file = out_path / f"{file.stem}_{file.suffix.replace('.', '')}.txt"
        output_file.write_text(text, encoding="utf-8")

        index.append({
            "source_file": str(file),
            "normalized_file": str(output_file),
            "size": len(text)
        })

    # Copy diagrams folder
    diagrams_src = raw_path / "diagrams"
    if diagrams_src.exists():
        diagrams_dest = out_path / "diagrams"
        shutil.copytree(diagrams_src, diagrams_dest, dirs_exist_ok=True)

        diagram_files = [str(p) for p in diagrams_dest.rglob("*") if p.is_file()]
        (out_path / "diagrams_manifest.json").write_text(
            json.dumps(diagram_files, indent=2),
            encoding="utf-8"
        )

    # Write architecture index
    (out_path / "architecture_index.json").write_text(
        json.dumps(index, indent=2),
        encoding="utf-8"
    )

    # Write failure log
    if failed_files:
        (out_path / "extraction_failures.json").write_text(
            json.dumps(failed_files, indent=2),
            encoding="utf-8"
        )

    print("Extraction complete.")
    print(f"Documents processed: {len(index)}")
    print(f"Failures: {len(failed_files)}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python extract_all_documents.py <raw_folder> <output_folder>")
        sys.exit(1)

    main(sys.argv[1], sys.argv[2])
